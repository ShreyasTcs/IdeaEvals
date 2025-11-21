import logging
from pathlib import Path
from typing import Dict, Optional
import platform
import subprocess
import tempfile
import json
import os
import logging
import base64
import pymupdf as fitz
from pathlib import Path
from typing import Dict
import cv2
from openai import AzureOpenAI
from PIL import Image
import io
from config.config import VIDEO_FRAME_EXTRACTION_INTERVAL_SECONDS
from pptx import Presentation
from docx import Document
import cv2

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

class VideoAnalyzer:
    """Analyze videos using Azure OpenAI Vision API"""
    
    def __init__(
        self,
        api_key: str,
        endpoint: str,
        deployment: str = "gpt-4.1",
        api_version: str = "2025-01-01-preview"
    ):
        self.client = AzureOpenAI(
            api_key=api_key,
            api_version=api_version,
            azure_endpoint=endpoint
        )
        self.deployment = deployment
    
    def analyze_video(self, video_path: str, interval_seconds: int = 10) -> Dict[str, str]:
        """
        Analyze video and return summary
        
        Args:
            video_path: Path to video file
            interval_seconds: Extract 1 frame every N seconds
        
        Returns:
            {
                "content": "Detailed video summary",
                "content_type": "Prototype" or "Text"
            }
        """
        video_path = Path(video_path)
        
        if not video_path.exists():
            raise FileNotFoundError(f"Video not found: {video_path}")
        
        logger.info(f"Analyzing video: {video_path.name}")
        
        # Step 1: Extract frames
        frames = self._extract_frames(video_path, interval_seconds)
        
        if not frames:
            return {
                'content': "Failed to extract frames from video",
                'content_type': 'Text'
            }
        
        logger.info(f"✓ Extracted {len(frames)} frames")
        
        # Step 2: Limit and compress frames
        frames_to_use = frames[:10]  # Max 10 frames for API
        frame_messages = self._prepare_frames_for_api(frames_to_use)
        
        if not frame_messages:
            return {
                'content': "Failed to prepare frames for analysis",
                'content_type': 'Text'
            }
        
        # Step 3: Send to Azure OpenAI Vision
        return self._analyze_with_vision(frame_messages, video_path.name)
    
    def _extract_frames(self, video_path: Path, interval_seconds: int) -> list:
        """Extract frames from video at regular intervals"""
        frames = []
        
        try:
            cap = cv2.VideoCapture(str(video_path))
            
            if not cap.isOpened():
                logger.error(f"Could not open video: {video_path}")
                return frames
            
            fps = cap.get(cv2.CAP_PROP_FPS)
            total_frames = int(cap.get(cv2.CAP_PROP_FRAME_COUNT))
            duration = total_frames / fps if fps > 0 else 0
            
            logger.info(f"Video info: {duration:.1f}s, {fps:.1f} FPS, {total_frames} frames")
            
            frame_interval = int(fps * VIDEO_FRAME_EXTRACTION_INTERVAL_SECONDS)
            frame_count = 0
            extracted_count = 0
            
            while True:
                ret, frame = cap.read()
                if not ret:
                    break
                
                if frame_count % frame_interval == 0:
                    # Encode frame as JPEG
                    success, buffer = cv2.imencode('.jpg', frame, [cv2.IMWRITE_JPEG_QUALITY, 85])
                    if success:
                        frames.append(buffer.tobytes())
                        extracted_count += 1
                        timestamp = frame_count / fps
                        logger.info(f"Frame {extracted_count} @ {timestamp:.1f}s")
                
                frame_count += 1
            
            cap.release()
            logger.info(f"✓ Total frames extracted: {len(frames)}")
            
        except Exception as e:
            logger.error(f"Frame extraction failed: {e}")
        
        return frames
    
    def _prepare_frames_for_api(self, frames: list) -> list:
        """Compress frames and convert to base64 for API"""
        frame_messages = []
        
        for i, frame_data in enumerate(frames):
            try:
                # Load image
                img = Image.open(io.BytesIO(frame_data))
                
                # Resize to reduce API payload size
                img.thumbnail((512, 512), Image.Resampling.LANCZOS)
                
                # Convert to JPEG with compression
                buffer = io.BytesIO()
                img.save(buffer, format='JPEG', quality=75)
                compressed_data = buffer.getvalue()
                
                # Encode to base64
                base64_frame = base64.b64encode(compressed_data).decode('utf-8')
                data_uri = f"data:image/jpeg;base64,{base64_frame}"
                
                frame_messages.append({
                    "type": "image_url",
                    "image_url": {
                        "url": data_uri,
                        "detail": "low"
                    }
                })
                
                logger.info(f"Prepared frame {i+1}: {len(compressed_data)} bytes")
                
            except Exception as e:
                logger.warning(f"Failed to prepare frame {i+1}: {e}")
                continue
        
        return frame_messages
    
    def _analyze_with_vision(self, frame_messages: list, video_name: str) -> Dict[str, str]:
        """Send frames to Azure OpenAI Vision for analysis"""
        
        prompt = f"""Analyze this video by examining {len(frame_messages)} frames extracted at regular intervals.

                    Video: {video_name}

                    Provide a comprehensive summary including:
                    1. **Main Topic/Purpose**: What is this video about?
                    2. **Key Points**: What are the main points shown across the frames?
                    3. **Visual Elements**: What UI, screens, diagrams, or demonstrations are shown?
                    4. **Technical Details**: Any technology, code, systems, or implementation shown?
                    5. **Workflow/Process**: What process or workflow is demonstrated?

                    Classification:
                    - **Prototype**: Shows working UI/demo, technical implementation, code, system architecture, working application
                    - **Text**: Only concepts, slides, presentations with no technical demonstration

                    Return ONLY valid JSON:
                    {{
                    "content": "Comprehensive detailed summary covering all 5 points above",
                    "content_type": "Prototype" or "Text"
                    }}"""
                            
        try:
            logger = logging.getLogger(__name__)
            logger.info(f"Analyzing {len(frame_messages)} frames with Azure OpenAI Vision...")
            
            response = self.client.chat.completions.create(
                model=self.deployment,
                messages=[
                    {
                        "role": "user",
                        "content": [{"type": "text", "text": prompt}] + frame_messages
                    }
                ],
                temperature=0.3,
                max_tokens=4000
            )
            
            result_text = response.choices[0].message.content
            
            # Clean markdown wrappers
            if result_text.startswith("```"):
                result_text = result_text.replace("```json", "").replace("```", "")

            
            result = json.loads(result_text)
            
            logger.info(f"✓ Analysis complete: {len(result['content'])} characters")
            logger.info(f"✓ Content type: {result['content_type']}")
            
            return result
            
        except Exception as e:
            logger.error(f"Vision analysis failed: {e}")
            return {
                'content': f"Video analysis failed: {str(e)}",
                'content_type': 'Text'
            }

class FileExtractor:
    """Process files using Azure OpenAI Vision API with fallback to text extraction"""
    
    def __init__(self):
        # Initialize Azure OpenAI client
        self.vision_client = AzureOpenAI(
            api_key=os.getenv("AZURE_OPENAI_API_KEY"),
            api_version="2025-01-01-preview",
            azure_endpoint=os.getenv("AZURE_OPENAI_ENDPOINT")
        )
        self.deployment = os.getenv("AZURE_OPENAI_MODEL", "gpt-4o")
    
    def extract_content(self, file_path: Path) -> Dict[str, str]:
        """
        Primary method: Analyze file with Azure OpenAI Vision API
        Fallback: Traditional text extraction if Vision API fails
        """
        logger.info(f"Processing file: {file_path.name}")
        
        # Ensure file exists
        if not file_path.exists():
            logger.error(f"File not found: {file_path}")
            return {"text": "", "content": "", "content_type": "not_found"}
        
        file_type = file_path.suffix.lower()
        result = None  # Initialize result variable
        
        # Try Vision API first if available
        if self.vision_client:
            try:
                if file_type in ['.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff']:
                    result = self._analyze_image_with_vision(file_path)
                    
                elif file_type == '.pdf' and fitz is not None:
                    result = self._analyze_pdf_with_vision(file_path)
                    
                elif file_type == '.pptx':
                    result = self._analyze_pptx_with_azure_di(file_path)
                    
                elif file_type == '.docx':
                    result = self._analyze_docx_with_vision(file_path)
                    
                elif file_type in ['.mp4', '.mov', '.avi']:
                    result = self._analyze_video_with_vision(file_path)
                
                # CRITICAL FIX: Add 'text' key if result was successfully obtained
                if result is not None:
                    result['text'] = result.get('content', '')
                    return result
                    
            except Exception as e:
                logger.warning(f"Vision API failed for {file_path.name}: {e}. Falling back to text extraction.")
        
        # Fallback to traditional extraction
        result = self._extract_text_traditional(file_path)
        result['text'] = result.get('content', '')
        return result
  
    def _analyze_image_with_vision(self, file_path: Path) -> Dict[str, str]:
        """Analyze image file directly with Azure OpenAI Vision"""
        logger.info(f"Analyzing image: {file_path.name}")
        
        try:
            # Encode image to base64
            with open(file_path, "rb") as image_file:
                base64_image = base64.b64encode(image_file.read()).decode('utf-8')
            
            ext = file_path.suffix.lower()
            mime_type = f"image/{'jpeg' if ext == '.jpg' else ext.strip('.')}"

            prompt = f"""Analyze this image and extract:
            1. Main content and purpose
            2. Text content (OCR)
            3. UI elements, diagrams, or technical visuals
            4. Code, architecture, or implementation details

            Is this a prototype/demo or just a concept slide?

            Return JSON:
            {{
            "content": "Detailed description",
            "content_type": "Prototype" or "Text"
            }}"""
            
            response = self.vision_client.chat.completions.create(
                model=self.deployment,
                messages=[{
                    "role": "user",
                    "content": [
                        {"type": "text", "text": prompt},
                        {"type": "image_url", "image_url": {"url": f"data:{mime_type};base64,{base64_image}"}}
                    ]
                }],
                temperature=0.3,
                max_tokens=2000
            )
            
            result_text = response.choices[0].message.content
            
            # Clean and parse JSON
            if result_text.startswith("```json"):
                result_text = result_text.replace("```json", "").replace("```", "")
            elif result_text.startswith("```"):
                result_text = result_text.replace("```", "")
            
            result = json.loads(result_text)
            
            logger.info(f"✓ Image analysis complete: {result['content_type']}")
            return result
            
        except Exception as e:
            logger.error(f"Image Vision analysis failed: {e}")
            raise  # Re-raise to trigger fallback
    
    def _analyze_document_with_vision(self, file_path: Path) -> Dict[str, str]:
        """Convert document pages to images and analyze with Vision API"""
        logger.info(f"Analyzing document: {file_path.name}")
        
        try:
            if file_path.suffix.lower() == '.pdf':
                return self._analyze_pdf_with_vision(file_path)
            elif file_path.suffix.lower() == '.docx':
                return self._analyze_docx_with_vision(file_path)
            elif file_path.suffix.lower() == '.pptx':
                return self._analyze_pptx_with_vision(file_path)
                
        except Exception as e:
            logger.error(f"Document Vision analysis failed: {e}")
            raise  # Re-raise to trigger fallback
    
    def _analyze_pdf_with_vision(self, file_path: Path) -> Dict[str, str]:
        """Convert PDF pages to images using PyMuPDF (no external dependencies)"""
        try:
            import fitz  # PyMuPDF
            
            logger.info(f"Rendering PDF pages with PyMuPDF: {file_path.name}")
            doc = fitz.open(file_path)
            
            frame_messages = []
            max_pages = min(len(doc), 5)  # Max 5 pages
            
            for page_num in range(max_pages):
                page = doc.load_page(page_num)
                
                # Render page to pixmap (image)
                zoom = 1.5  # zoom factor for better quality
                mat = fitz.Matrix(zoom, zoom)
                pix = page.get_pixmap(matrix=mat)
                
                # Convert to PIL Image
                img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                
                # Convert to base64
                buffer = io.BytesIO()
                img.save(buffer, format='JPEG', quality=85)
                base64_image = base64.b64encode(buffer.getvalue()).decode('utf-8')
                
                frame_messages.append({
                    "type": "image_url",
                    "image_url": {
                        "url": f"data:image/jpeg;base64,{base64_image}",
                        "detail": "low"
                    }
                })
                logger.info(f"✓ Rendered PDF page {page_num+1}/{max_pages}")
            
            doc.close()
            
            return self._analyze_with_vision_api(
                frame_messages, 
                file_path.name,
                doc_type="PDF document"
            )
            
        except ImportError:
            logger.error("PyMuPDF not installed. Install with: pip install PyMuPDF")
            raise
        except Exception as e:
            logger.error(f"PDF Vision analysis failed: {e}")
            raise

    def _analyze_docx_with_vision(self, file_path: Path) -> Dict[str, str]:
        """Analyze DOCX by extracting text and sending as structured data"""
        try:
            doc = Document(file_path)
            full_text = []
            
            for para in doc.paragraphs:
                if para.text.strip():
                    full_text.append(para.text)
            
            for table in doc.tables:
                for row in table.rows:
                    row_text = [cell.text.strip() for cell in row.cells]
                    full_text.append(" | ".join(row_text))
            
            content = "\n".join(full_text)
            
            prompt = f"""Analyze this Word document content:

                        {content[:4000]}  # Truncate to avoid token limits

                        Determine:
                        1. Main topic and purpose
                        2. Technical content, code, or implementation details
                        3. Is this a concept document or does it describe a working prototype?

                        Return JSON:
                        {{
                        "content": "Detailed analysis",
                        "content_type": "Prototype" or "Text"
                        }}"""
            
            response = self.vision_client.chat.completions.create(
                model=self.deployment,
                messages=[{"role": "user", "content": prompt}],
                temperature=0.3,
                max_tokens=1000
            )
            
            result_text = response.choices[0].message.content
            
            # Clean and parse JSON
            if result_text.startswith("```json"):
                result_text = result_text.replace("```json", "").replace("```", "")
            elif result_text.startswith("```"):
                result_text = result_text.replace("```", "")
            
            result = json.loads(result_text)
            result['content'] = content[:2000]  # Include actual content
            
            logger.info(f"✓ DOCX analysis complete: {result['content_type']}")
            return result
            
        except Exception as e:
            logger.error(f"DOCX Vision analysis failed: {e}")
            raise
    
    def _analyze_pptx_with_azure_di(self, file_path: Path) -> Dict[str, str]:
        """
        Extract text from PPTX using Azure Document Intelligence
        Fixed for latest SDK version
        """
        from azure.ai.documentintelligence import DocumentIntelligenceClient
        from azure.ai.documentintelligence.models import AnalyzeDocumentRequest
        from azure.core.credentials import AzureKeyCredential
        
        logger.info(f"Extracting text from PPTX with Azure DI: {file_path.name}")
        
        try:
            # Initialize client
            client = DocumentIntelligenceClient(
                endpoint=os.getenv("AZURE_OPENAI_ENDPOINT"),
                credential=AzureKeyCredential(os.getenv("AZURE_OPENAI_API_KEY"))
            )
            
            # Read file as bytes
            with open(file_path, "rb") as f:
                file_bytes = f.read()
            
            # Analyze document with correct parameters
            poller = client.begin_analyze_document(
                model_id="prebuilt-read",
                body=file_bytes,  # Pass bytes directly as 'body'
                content_type="application/octet-stream"
            )
            
            # Wait for result
            result = poller.result()
            
            # Extract text safely
            all_text = []
            
            # Try pages -> lines
            if hasattr(result, 'pages') and result.pages:
                for page in result.pages:
                    if hasattr(page, 'lines') and page.lines:
                        for line in page.lines:
                            if line.content:
                                all_text.append(line.content)
            
            # Try paragraphs as fallback
            if not all_text and hasattr(result, 'paragraphs') and result.paragraphs:
                for paragraph in result.paragraphs:
                    if paragraph.content:
                        all_text.append(paragraph.content)
            
            content = "\n".join(all_text)
            
            if not content:
                raise ValueError("No content extracted from document")
            
            logger.info(f"✓ Extracted {len(content)} characters from Azure DI")
            
            return {
                'content': content,
                'text': content,
                'content_type': 'Text'
            }
            
        except Exception as e:
            logger.error(f"Azure Document Intelligence failed: {e}")
            raise

    def _analyze_video_with_vision(self, file_path: Path) -> Dict[str, str]:
        """Analyze video using frame extraction and Vision API"""
        try:            
            analyzer = VideoAnalyzer(
                api_key=os.getenv("AZURE_OPENAI_API_KEY"),
                endpoint=os.getenv("AZURE_OPENAI_ENDPOINT"),
                deployment=os.getenv("AZURE_OPENAI_MODEL", "gpt-4o")
            )
            
            return analyzer.analyze_video(
                video_path=file_path,
                interval_seconds=VIDEO_FRAME_EXTRACTION_INTERVAL_SECONDS
            )
            
        except Exception as e:
            logger.error(f"Video Vision analysis failed: {e}")
            raise
    
    def _analyze_with_vision_api(self, frame_messages: list, file_name: str, doc_type: str = "document") -> Dict[str, str]:
        """Generic method to analyze images with Vision API"""
        
        prompt = f"""Analyze this {doc_type}: {file_name}

                    You have {len(frame_messages)} pages/slides/frames. Extract:
                    1. Main content and purpose
                    2. Technical details, code, implementation
                    3. Visual elements, UI, diagrams
                    4. Overall classification: Prototype or Text

                    Return JSON:
                    {{
                    "content": "Comprehensive analysis",
                    "content_type": "Prototype" or "Text"
                    }}"""
        
        try:
            response = self.vision_client.chat.completions.create(
                model=self.deployment,
                messages=[{
                    "role": "user",
                    "content": [{"type": "text", "text": prompt}] + frame_messages
                }],
                temperature=0.3,
                max_tokens=4000
            )
            
            result_text = response.choices[0].message.content
            
            # Clean and parse JSON
            if result_text.startswith("```json"):
                result_text = result_text.replace("```json", "").replace("```", "")
            elif result_text.startswith("```"):
                result_text = result_text.replace("```", "")
            
            result = json.loads(result_text)
            logger.info(f"✓ {doc_type} analysis complete: {result['content_type']}")
            return result
            
        except Exception as e:
            logger.error(f"Vision API analysis failed: {e}")
            raise
    
    def _extract_text_traditional(self, file_path: Path) -> Dict[str, str]:
        """Fallback: Traditional text extraction for all file types"""
        logger.info(f"Using traditional extraction for: {file_path.name}")
        
        file_type = file_path.suffix.lower()
        
        try:
            if file_type == '.pdf':
                # Try PyMuPDF first
                try:
                    import pymupdf as fitz
                    doc = fitz.open(file_path)
                    text = ""
                    for page in doc:
                        text += page.get_text()
                    return {"content": text, "content_type": "Text"}
                except:
                    return {"content": f"PDF file: {file_path.name}", "content_type": "Text"}
            
            elif file_type == '.pptx':
                return {"content": self._extract_text_from_pptx(file_path), "content_type": "Text"}
            
            elif file_type == '.docx':
                return {"content": self._extract_text_from_docx(file_path), "content_type": "Text"}
            
            elif file_type in ['.jpg', '.jpeg', '.png']:
                return {"content": f"Image file: {file_path.name}", "content_type": "Text"}
            
            elif file_type in ['.mp4', '.mov', '.avi']:
                return {"content": self._extract_info_from_video(file_path), "content_type": "Text"}
            
            else:
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    return {"content": f.read(), "content_type": "Text"}
                    
        except Exception as e:
            logger.error(f"Traditional extraction failed for {file_path.name}: {e}")
            return {"content": f"Error extracting {file_path.name}", "content_type": "Text"}
    
    def _extract_text_from_pptx(self, file_path: Path) -> str:
        """Extracts text from a PPTX file."""
        try:
            prs = Presentation(file_path)
            text_runs = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_runs.append(shape.text)
            return "\n".join(text_runs)
        except Exception as e:
            logger.error(f"Error reading PPTX file {file_path}: {e}")
            return f"Could not extract text from {file_path.name}."
    
    def _extract_text_from_docx(self, file_path: Path) -> str:
        """Extracts text from a DOCX file."""
        try:
            doc = Document(file_path)
            return "\n".join([p.text for p in doc.paragraphs])
        except Exception as e:
            logger.error(f"Error reading DOCX file {file_path}: {e}")
            return f"Could not extract text from {file_path.name}."
    
    def _extract_info_from_video(self, file_path: Path) -> str:
        """Extracts basic information from a video file."""
        try:
            import cv2
            cap = cv2.VideoCapture(str(file_path))
            if not cap.isOpened():
                raise IOError("Could not open video file.")
            
            fps = cap.get(cv2.CAP_PROP_FPS)
            frame_count = int(cap.get(cv2.CAP_PROP_FRAME_COUNT))
            duration = frame_count / fps if fps > 0 else 0
            width = int(cap.get(cv2.CAP_PROP_FRAME_WIDTH))
            height = int(cap.get(cv2.CAP_PROP_FRAME_HEIGHT))
            
            cap.release()
            return f"Video file: {file_path.name}. Duration: {duration:.2f}s, Resolution: {width}x{height}"
        except Exception as e:
            logger.error(f"Error processing video file {file_path}: {e}")
            return f"Error processing video file {file_path.name}."
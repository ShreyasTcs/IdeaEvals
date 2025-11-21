import logging
from pathlib import Path
from typing import List
import cv2
from PIL import Image
import io
import base64
from pptx import Presentation
from docx import Document

logger = logging.getLogger(__name__)

def extract_frames(video_path: Path, interval_seconds: int) -> list:
    """Extract frames from video at regular intervals"""
    frames = []
    
    try:
        cap = cv2.VideoCapture(str(video_path))
        
        if not cap.isOpened():
            logger.error(f"Could not open video: {video_path}")
            return frames
        
        fps = cap.get(cv2.CAP_PROP_FPS)
        total_frames = int(cap.get(cv2.CAP_PROP_FRAME_COUNT))
        duration = total_frames / fps if fps > 0 else 0
        
        logger.info(f"Video info: {duration:.1f}s, {fps:.1f} FPS, {total_frames} frames")
        
        frame_interval = int(fps * interval_seconds)
        frame_count = 0
        extracted_count = 0
        
        while True:
            ret, frame = cap.read()
            if not ret:
                break
            
            if frame_count % frame_interval == 0:
                # Encode frame as JPEG
                success, buffer = cv2.imencode('.jpg', frame, [cv2.IMWRITE_JPEG_QUALITY, 85])
                if success:
                    frames.append(buffer.tobytes())
                    extracted_count += 1
                    timestamp = frame_count / fps
                    logger.info(f"Frame {extracted_count} @ {timestamp:.1f}s")
            
            frame_count += 1
        
        cap.release()
        logger.info(f"✓ Total frames extracted: {len(frames)}")
        
    except Exception as e:
        logger.error(f"Frame extraction failed: {e}")
    
    return frames

def prepare_frames_for_api(frames: list) -> list:
    """Compress frames and convert to base64 for API"""
    frame_messages = []
    
    for i, frame_data in enumerate(frames):
        try:
            # Load image
            img = Image.open(io.BytesIO(frame_data))
            
            # Resize to reduce API payload size
            img.thumbnail((512, 512), Image.Resampling.LANCZOS)
            
            # Convert to JPEG with compression
            buffer = io.BytesIO()
            img.save(buffer, format='JPEG', quality=75)
            compressed_data = buffer.getvalue()
            
            # Encode to base64
            base64_frame = base64.b64encode(compressed_data).decode('utf-8')
            data_uri = f"data:image/jpeg;base64,{base64_frame}"
            
            frame_messages.append({
                "type": "image_url",
                "image_url": {
                    "url": data_uri,
                    "detail": "low"
                }
            })
            
            logger.info(f"Prepared frame {i+1}: {len(compressed_data)} bytes")
            
        except Exception as e:
            logger.warning(f"Failed to prepare frame {i+1}: {e}")
            continue
    
    return frame_messages

def extract_text_from_pptx(file_path: Path) -> str:
    """Extracts text from a PPTX file."""
    try:
        prs = Presentation(file_path)
        text_runs = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_runs.append(shape.text)
        return "\n".join(text_runs)
    except Exception as e:
        logger.error(f"Error reading PPTX file {file_path}: {e}")
        return f"Could not extract text from {file_path.name}."

def extract_text_from_docx(file_path: Path) -> str:
    """Extracts text from a DOCX file."""
    try:
        doc = Document(file_path)
        return "\n".join([p.text for p in doc.paragraphs])
    except Exception as e:
        logger.error(f"Error reading DOCX file {file_path}: {e}")
        return f"Could not extract text from {file_path.name}."

def extract_info_from_video(file_path: Path) -> str:
    """Extracts basic information from a video file."""
    try:
        import cv2
        cap = cv2.VideoCapture(str(file_path))
        if not cap.isOpened():
            raise IOError("Could not open video file.")
        
        fps = cap.get(cv2.CAP_PROP_FPS)
        frame_count = int(cap.get(cv2.CAP_PROP_FRAME_COUNT))
        duration = frame_count / fps if fps > 0 else 0
        width = int(cap.get(cv2.CAP_PROP_FRAME_WIDTH))
        height = int(cap.get(cv2.CAP_PROP_FRAME_HEIGHT))
        
        cap.release()
        return f"Video file: {file_path.name}. Duration: {duration:.2f}s, Resolution: {width}x{height}"
    except Exception as e:
        logger.error(f"Error processing video file {file_path}: {e}")
        return f"Error processing video file {file_path.name}."
